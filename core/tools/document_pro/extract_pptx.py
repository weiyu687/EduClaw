"""
pptx 内容 (skill: Document Pro)

Author: Gongmin Wei
Date: 2026-04-07
"""
import json
from pptx import Presentation


def extract_pptx(pptx_path: str) -> str:
    """
    从 PPTX 文件中提取内容 (文本 + 表格)
    :param pptx_path: PPTX 文件路径
    :return: 序列化的 JSON 字符串，包含每页幻灯片的文本和表格
    """
    content = {
        "total_slides": 0,
        "slides": []
    }

    try:
        prs = Presentation(pptx_path)
        content["total_slides"] = len(prs.slides)

        for i, slide in enumerate(prs.slides, 1):
            slide_content = {
                "slide_num": i,
                "text": [],
                "tables": []
            }

            for shape in slide.shapes:
                # 提取文本内容
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_content["text"].append(text)

                # 提取表格内容并转换为 Markdown 格式
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        rows.append([cell.text.strip() for cell in row.cells])

                    if rows:
                        md_table = []
                        for j, row in enumerate(rows):
                            line = "| " + " | ".join(row) + " |"
                            md_table.append(line)
                            if j == 0:
                                md_table.append("|" + "|".join(["---"] * len(row)) + "|")
                        slide_content["tables"].append("\n".join(md_table))

            content["slides"].append(slide_content)

    except Exception as e:
        return json.dumps({"error": f"PPTX 文件读取失败: {str(e)}"}, ensure_ascii=False, indent=2)

    return json.dumps(content, ensure_ascii=False, indent=2)
